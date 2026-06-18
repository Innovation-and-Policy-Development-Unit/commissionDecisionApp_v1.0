"""
Compliance module API.

Everything runs inside SCDMS — there is no external system.

  * Compliance staff (manager / senior / principal) create and manage cases, and
    triage complaints from the Complaints Register.
  * Ministry staff (Head of Agency / HR) lodge complaints (write-only) and see only
    their own complaint + its coarse status — never the resulting case.

Phase 4 adds the full visibility-scoping test matrix; this phase implements the
queryset scoping and per-action permissions.
"""

from __future__ import annotations

import io
from datetime import date

from django.http import HttpResponse
from rest_framework import permissions, status, viewsets
from rest_framework.decorators import action
from rest_framework.exceptions import PermissionDenied, ValidationError
from rest_framework.response import Response

from .compliance_actions import (
    approve_compliance_case,
    create_compliance_case,
    return_compliance_case,
    submit_compliance_case,
)
from .compliance_models import Complaint, ComplaintStatus, ComplianceCase, OffenceType
from .compliance_scoping import (
    MINISTRY_LODGE_ROLES,
    complaint_queryset,
    compliance_case_queryset,
    user_can_manage_senior_cases,
    user_can_record_decision,
    user_can_view_compliance,
    user_is_compliance_manager,
    user_is_dg_director,
)
from .compliance_serializers import (
    CaseNoteSerializer,
    ComplaintLodgeSerializer,
    ComplaintSerializer,
    ComplianceCaseCreateSerializer,
    ComplianceCaseDecisionSerializer,
    ComplianceCaseDetailSerializer,
    ComplianceCaseListSerializer,
    GrievanceMediatorSerializer,
    LitigationRecordSerializer,
)


def _is_compliance_staff(user) -> bool:
    return user_can_view_compliance(user)


def _user_role(user) -> str | None:
    from .profile_utils import ensure_psc_profile

    try:
        return ensure_psc_profile(user).role
    except Exception:
        return None


class IsComplianceStaff(permissions.BasePermission):
    """Anyone with any compliance visibility may read; writes are action-gated."""

    message = "Only Compliance unit staff may access compliance cases."

    def has_permission(self, request, view):
        if not request.user or not request.user.is_authenticated:
            return False
        return user_can_view_compliance(request.user)


# Actions that require full compliance-staff write access (not just view).
_WRITE_ACTIONS = {"create", "submit", "approve", "return_for_changes", "notes", "litigation", "update_litigation"}

def _require_write_access(user):
    """Raise PermissionDenied if the user is a read-only compliance role."""
    from .compliance_scoping import COMPLIANCE_STAFF_ROLES
    from .compliance_scoping import _role as get_role
    from .models import Role
    read_only = {Role.DG_DIRECTOR, Role.COMMISSION_MEMBER, Role.PANEL_MEMBER}
    if not user_can_view_compliance(user):
        raise PermissionDenied("No access.")
    role = get_role(user)
    if role in read_only:
        raise PermissionDenied("Your role has read-only access to compliance cases.")


class ComplianceCaseViewSet(viewsets.ModelViewSet):
    """Create, list, retrieve, and run the approval flow for compliance cases."""

    permission_classes = [permissions.IsAuthenticated, IsComplianceStaff]
    http_method_names = ["get", "post", "head", "options"]

    def get_queryset(self):
        base = (
            ComplianceCase.objects
            .select_related("submission")
            .prefetch_related("stages", "litigation_records", "case_notes", "decisions")
            .order_by("-created_at")
        )
        return compliance_case_queryset(self.request.user, base)

    def get_serializer_class(self):
        if self.action == "create":
            return ComplianceCaseCreateSerializer
        if self.action == "retrieve":
            return ComplianceCaseDetailSerializer
        return ComplianceCaseListSerializer

    def create(self, request, *args, **kwargs):
        _require_write_access(request.user)
        ser = ComplianceCaseCreateSerializer(data=request.data)
        ser.is_valid(raise_exception=True)
        data = dict(ser.validated_data)

        complaint = None
        complaint_id = data.pop("complaint_id", None)
        if complaint_id:
            complaint = Complaint.objects.filter(pk=complaint_id).first()
            if not complaint:
                raise ValidationError({"complaint_id": "Complaint not found."})

        try:
            case = create_compliance_case(creator=request.user, complaint=complaint, **data)
        except ValueError as exc:
            raise ValidationError(str(exc)) from exc

        out = ComplianceCaseDetailSerializer(case, context={"request": request})
        return Response(out.data, status=status.HTTP_201_CREATED)

    @action(detail=True, methods=["post"])
    def submit(self, request, pk=None):
        """Submit the case: Senior/Principal → Manager approval; Manager → Secretary."""
        _require_write_access(request.user)
        case = self.get_object()
        submit_compliance_case(case, request.user, _user_role(request.user))
        return Response(ComplianceCaseDetailSerializer(case, context={"request": request}).data)

    @action(detail=True, methods=["post"])
    def approve(self, request, pk=None):
        """Compliance Manager approves → forwarded to Secretary review."""
        if not user_is_compliance_manager(request.user):
            raise PermissionDenied("Only the Compliance Manager may approve a case.")
        case = self.get_object()
        approve_compliance_case(case, request.user)
        return Response(ComplianceCaseDetailSerializer(case, context={"request": request}).data)

    @action(detail=True, methods=["post"], url_path="return")
    def return_for_changes(self, request, pk=None):
        """Compliance Manager returns the case to the originating officer."""
        if not user_is_compliance_manager(request.user):
            raise PermissionDenied("Only the Compliance Manager may return a case.")
        case = self.get_object()
        return_compliance_case(case, request.user, reason=request.data.get("reason", ""))
        return Response(ComplianceCaseDetailSerializer(case, context={"request": request}).data)

    @action(detail=True, methods=["post"])
    def notes(self, request, pk=None):
        """Add a case note."""
        _require_write_access(request.user)
        case = self.get_object()
        ser = CaseNoteSerializer(data=request.data)
        ser.is_valid(raise_exception=True)
        ser.save(case=case, author=request.user)
        return Response(ser.data, status=status.HTTP_201_CREATED)

    @action(detail=True, methods=["post"])
    def litigation(self, request, pk=None):
        """Add a litigation record (FR-13 cost tracking)."""
        _require_write_access(request.user)
        case = self.get_object()
        ser = LitigationRecordSerializer(data=request.data)
        ser.is_valid(raise_exception=True)
        ser.save(case=case)
        return Response(
            ComplianceCaseDetailSerializer(case, context={"request": request}).data,
            status=status.HTTP_201_CREATED,
        )

    @action(detail=True, methods=["patch"], url_path="litigation/(?P<lit_id>[0-9]+)")
    def update_litigation(self, request, pk=None, lit_id=None):
        """Update status, costs, court date, or notes on an existing litigation record."""
        from .compliance_models import LitigationRecord as LR
        case = self.get_object()
        rec = case.litigation_records.filter(pk=lit_id).first()
        if not rec:
            raise ValidationError({"detail": "Litigation record not found."})
        ser = LitigationRecordSerializer(rec, data=request.data, partial=True)
        ser.is_valid(raise_exception=True)
        ser.save()
        return Response(
            ComplianceCaseDetailSerializer(case, context={"request": request}).data
        )

    @action(detail=True, methods=["post", "put"], url_path="mediator")
    def mediator(self, request, pk=None):
        """Create or replace the mediator appointment for a grievance case (FR-11)."""
        from .compliance_models import GrievanceMediatorAppointment

        case = self.get_object()
        if case.case_family != "grievance":
            raise ValidationError({"detail": "Mediator appointments are only for grievance cases."})
        if not user_can_manage_senior_cases(request.user):
            raise PermissionDenied("Only the Compliance Manager or Secretary OPSC may appoint mediators.")

        existing = getattr(case, "mediator_appointment", None)
        ser = GrievanceMediatorSerializer(
            existing, data=request.data,
            partial=(request.method == "PUT"),
        )
        ser.is_valid(raise_exception=True)
        ser.save(case=case, appointed_by=request.user)
        return Response(
            ComplianceCaseDetailSerializer(case, context={"request": request}).data,
            status=status.HTTP_201_CREATED if not existing else status.HTTP_200_OK,
        )

    @action(detail=True, methods=["get", "post", "patch"], url_path="investigation")
    def investigation(self, request, pk=None):
        """Get, create, or update the structured investigation for a case
        (panel, terms of reference, findings, recommendation, report)."""
        from .compliance_models import Investigation
        from .compliance_serializers import InvestigationSerializer

        case = self.get_object()
        existing = getattr(case, "investigation", None)

        if request.method == "GET":
            if not existing:
                return Response({"detail": "No investigation recorded."}, status=404)
            return Response(InvestigationSerializer(existing).data)

        _require_write_access(request.user)
        ser = InvestigationSerializer(
            existing, data=request.data, partial=(request.method == "PATCH"),
        )
        ser.is_valid(raise_exception=True)
        if existing:
            ser.save()
        else:
            ser.save(case=case, created_by=request.user)
        return Response(
            ComplianceCaseDetailSerializer(case, context={"request": request}).data,
            status=status.HTTP_201_CREATED if not existing else status.HTTP_200_OK,
        )

    @action(detail=True, methods=["post"], url_path="suspensions")
    def add_suspension(self, request, pk=None):
        """Record a suspension and its financial implication (half/full/no-pay)."""
        from .compliance_serializers import SuspensionRecordSerializer

        _require_write_access(request.user)
        case = self.get_object()
        ser = SuspensionRecordSerializer(data=request.data)
        ser.is_valid(raise_exception=True)
        ser.save(case=case, created_by=request.user)
        return Response(
            ComplianceCaseDetailSerializer(case, context={"request": request}).data,
            status=status.HTTP_201_CREATED,
        )

    @action(detail=True, methods=["patch"], url_path="suspensions/(?P<sus_id>[0-9]+)")
    def update_suspension(self, request, pk=None, sus_id=None):
        """Update suspension fields (amounts, dates, notes)."""
        from .compliance_serializers import SuspensionRecordSerializer

        _require_write_access(request.user)
        case = self.get_object()
        rec = case.suspensions.filter(pk=sus_id).first()
        if not rec:
            raise ValidationError({"detail": "Suspension record not found."})
        ser = SuspensionRecordSerializer(rec, data=request.data, partial=True)
        ser.is_valid(raise_exception=True)
        ser.save()
        return Response(ComplianceCaseDetailSerializer(case, context={"request": request}).data)

    @action(detail=True, methods=["post"], url_path="suspensions/(?P<sus_id>[0-9]+)/assess")
    def assess_suspension_reimbursement(self, request, pk=None, sus_id=None):
        """OPSC assessment of whether withheld salary is reimbursed on reinstatement.
        Body: { reimbursement_status: reimburse|forfeit, reimbursed_amount? }."""
        from django.utils import timezone as _tz
        from .compliance_models import SuspensionReimbursement

        if not user_can_manage_senior_cases(request.user):
            raise PermissionDenied("Only the Compliance Manager or Secretary OPSC may assess reimbursement.")
        case = self.get_object()
        rec = case.suspensions.filter(pk=sus_id).first()
        if not rec:
            raise ValidationError({"detail": "Suspension record not found."})
        decision = request.data.get("reimbursement_status")
        if decision not in (SuspensionReimbursement.REIMBURSE, SuspensionReimbursement.FORFEIT):
            raise ValidationError({"reimbursement_status": "Must be 'reimburse' or 'forfeit'."})
        rec.reimbursement_status = decision
        rec.opsc_assessed_by = request.user
        rec.opsc_assessed_at = _tz.now()
        if decision == SuspensionReimbursement.REIMBURSE and request.data.get("reimbursed_amount") is not None:
            rec.reimbursed_amount = request.data["reimbursed_amount"]
        rec.save()
        return Response(ComplianceCaseDetailSerializer(case, context={"request": request}).data)

    @action(detail=True, methods=["get", "post"], url_path="documents")
    def documents(self, request, pk=None):
        """List or upload documents/evidence for a case, stored against its
        submission. ``doc_type`` tags the upload by workflow form (SMDR, warning,
        response, investigation report, evidence, outcome letter)."""
        from .models import SubmissionDocument
        from .serializers import SubmissionDocumentSerializer

        case = self.get_object()
        submission = case.submission

        if request.method == "GET":
            docs = SubmissionDocument.objects.filter(submission=submission)
            return Response(SubmissionDocumentSerializer(docs, many=True, context={"request": request}).data)

        _require_write_access(request.user)
        f = request.FILES.get("file")
        if not f:
            return Response({"detail": "No file provided."}, status=status.HTTP_400_BAD_REQUEST)
        if f.size > 20 * 1024 * 1024:
            return Response({"detail": f"File '{f.name}' exceeds the 20 MB limit."}, status=status.HTTP_400_BAD_REQUEST)
        doc = SubmissionDocument.objects.create(
            submission=submission,
            file=f,
            original_name=f.name,
            description=(request.data.get("doc_type") or request.data.get("description") or ""),
            uploaded_by=request.user,
        )
        return Response(
            SubmissionDocumentSerializer(doc, context={"request": request}).data,
            status=status.HTTP_201_CREATED,
        )

    @action(detail=True, methods=["post"], url_path="documents/(?P<doc_id>[0-9]+)/remove")
    def remove_document(self, request, pk=None, doc_id=None):
        """Soft-delete (archive) a case document."""
        from django.utils import timezone as _tz
        from .models import SubmissionDocument

        _require_write_access(request.user)
        case = self.get_object()
        doc = SubmissionDocument.objects.filter(submission=case.submission, pk=doc_id).first()
        if not doc:
            raise ValidationError({"detail": "Document not found."})
        doc.archived_at = _tz.now()
        doc.archived_by = request.user
        doc.save(update_fields=["archived_at", "archived_by"])
        return Response({"detail": "Document removed."})

    @action(detail=True, methods=["get"], url_path="documents/(?P<doc_id>[0-9]+)/download")
    def download_document(self, request, pk=None, doc_id=None):
        """Stream a case document to compliance staff who can view the case."""
        import mimetypes

        from django.http import FileResponse
        from django.shortcuts import get_object_or_404

        from .models import SubmissionDocument

        case = self.get_object()
        doc = get_object_or_404(SubmissionDocument, id=doc_id, submission=case.submission)
        content_type, _ = mimetypes.guess_type(doc.original_name)
        resp = FileResponse(doc.file.open("rb"), content_type=content_type or "application/octet-stream")
        resp["Content-Disposition"] = f'inline; filename="{doc.original_name}"'
        return resp

    @action(detail=False, methods=["get"], url_path="export-pptx")
    def export_pptx(self, request):
        """Generate and return a PPTX summary of all compliance cases."""
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        cases = list(self.get_queryset().order_by("created_at"))

        NAVY  = RGBColor(0x1E, 0x3A, 0x5F)
        GOLD  = RGBColor(0xC9, 0xA8, 0x4C)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        LGRAY = RGBColor(0xF0, 0xF4, 0xF8)
        RED   = RGBColor(0xDC, 0x26, 0x26)
        GREEN = RGBColor(0x16, 0xA3, 0x4A)
        DKTXT = RGBColor(0x1E, 0x29, 0x3B)

        def rgb_fill(shape, color):
            from pptx.oxml.ns import qn
            from lxml import etree
            sp = shape.fill
            sp.solid()
            sp.fore_color.rgb = color

        prs = Presentation()
        prs.slide_width  = Inches(10)
        prs.slide_height = Inches(5.625)
        blank = prs.slide_layouts[6]  # blank layout

        def add_rect(slide, l, t, w, h, fill_color, line_color=None):
            from pptx.util import Inches
            shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
            shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
            if line_color:
                shape.line.color.rgb = line_color
                shape.line.width = Pt(0.5)
            else:
                shape.line.fill.background()
            return shape

        def add_text(slide, text, l, t, w, h, size=11, bold=False, color=None, align="left", valign="middle"):
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
            run = p.add_run()
            run.text = str(text)
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color or DKTXT
            return txBox

        def header(slide, title):
            add_rect(slide, 0, 0, 10, 0.72, NAVY)
            add_text(slide, title, 0.4, 0.08, 9.2, 0.56, size=18, bold=True, color=WHITE)

        def slide_footer(slide):
            add_rect(slide, 0, 5.42, 10, 0.2, GOLD)
            add_text(slide, "Public Service Commission — Compliance & Integrity Unit  |  Confidential",
                     0, 5.42, 10, 0.2, size=7, color=NAVY, align="center")

        # ── Slide 1: Title ──────────────────────────────────────────────────────
        s1 = prs.slides.add_slide(blank)
        add_rect(s1, 0, 0, 10, 5.625, NAVY)
        add_rect(s1, 0, 4.2, 10, 0.15, GOLD)
        add_text(s1, "SUMMARY REPORT", 1, 1.0, 8, 0.5, size=13, bold=True, color=GOLD, align="center")
        add_text(s1, "Pending Disciplinary Matters", 1, 1.55, 8, 0.9, size=32, bold=True, color=WHITE, align="center")
        add_text(s1, "Public Service Commission — Compliance & Integrity Unit", 1, 2.6, 8, 0.4, size=13, color=RGBColor(0xA8,0xC4,0xE0), align="center")
        add_text(s1, f"Generated: {date.today().strftime('%d %B %Y')}", 1, 3.1, 8, 0.35, size=12, color=GOLD, align="center")
        add_text(s1, f"Total Cases: {len(cases)}", 1, 3.55, 8, 0.35, size=12, color=WHITE, align="center")

        # ── Slide 2: Stats snapshot ─────────────────────────────────────────────
        s2 = prs.slides.add_slide(blank)
        add_rect(s2, 0, 0, 10, 5.625, LGRAY)
        header(s2, "Compliance Cases — Snapshot")

        open_cases   = [c for c in cases if c.status not in ("closed", "archived")]
        closed_cases = [c for c in cases if c.status == "closed"]
        overdue      = [c for c in cases if (getattr(c, "sla_summary", None) or {}).get("overdue", 0) > 0]
        senior       = [c for c in cases if getattr(c, "is_senior_executive", False)]

        for i, (val, lbl, col) in enumerate([
            (len(cases),       "Total Cases",  NAVY),
            (len(open_cases),  "Open Cases",   NAVY),
            (len(overdue),     "Overdue SLA",  RED),
            (len(closed_cases),"Closed",       GREEN),
        ]):
            xx = 0.5 + i * 2.3
            add_rect(s2, xx, 0.9, 2.0, 1.3, WHITE, LGRAY)
            add_text(s2, str(val), xx, 1.0, 2.0, 0.7, size=36, bold=True, color=col, align="center")
            add_text(s2, lbl,      xx, 1.75, 2.0, 0.35, size=10, color=RGBColor(0x64,0x74,0x8B), align="center")

        # Family breakdown
        families = {}
        for c in cases:
            fam = getattr(c, "case_family_display", None) or getattr(c, "case_family", "Unknown")
            families[fam] = families.get(fam, 0) + 1

        add_text(s2, "CASELOAD BY FAMILY", 0.5, 2.38, 9, 0.28, size=9, bold=True, color=GOLD)
        add_rect(s2, 0.5, 2.68, 9, 0.3, NAVY)
        add_text(s2, "Family", 0.55, 2.68, 5, 0.3, size=9, bold=True, color=WHITE)
        add_text(s2, "Cases", 8.8, 2.68, 0.9, 0.3, size=9, bold=True, color=WHITE, align="center")

        for row_i, (fam, cnt) in enumerate(sorted(families.items(), key=lambda x: -x[1])[:7]):
            yy = 2.98 + row_i * 0.3
            bg = WHITE if row_i % 2 == 0 else LGRAY
            add_rect(s2, 0.5, yy, 9, 0.3, bg, RGBColor(0xCB,0xD5,0xE1))
            add_text(s2, fam,      0.55, yy, 7.8, 0.3, size=9,  color=DKTXT)
            add_text(s2, str(cnt), 8.8,  yy, 0.9, 0.3, size=9, bold=True, color=NAVY, align="center")

        slide_footer(s2)

        # ── Slide 3: Active Cases register ──────────────────────────────────────
        s3 = prs.slides.add_slide(blank)
        add_rect(s3, 0, 0, 10, 5.625, LGRAY)
        header(s3, f"Active Disciplinary Cases Register  [{len(open_cases)} open]")

        add_rect(s3, 0.3, 0.85, 9.4, 0.3, NAVY)
        for hdr, lx, lw in [("Officer / Subject", 0.35, 2.8), ("Family / Type", 3.2, 2.6), ("Status", 5.85, 2.2), ("Ref", 8.1, 1.5)]:
            add_text(s3, hdr, lx, 0.85, lw, 0.3, size=8.5, bold=True, color=WHITE)

        row_h = 0.34
        for ri, c in enumerate(open_cases[:13]):
            yy = 1.15 + ri * row_h
            bg = WHITE if ri % 2 == 0 else LGRAY
            add_rect(s3, 0.3, yy, 9.4, row_h, bg, RGBColor(0xCB,0xD5,0xE1))
            name   = getattr(c, "subject_name", str(c))
            family = getattr(c, "case_family_display", getattr(c, "case_family", "—"))
            cstatus = getattr(c, "get_status_display", lambda: getattr(c, "status", "—"))()
            ref    = getattr(c, "reference_number", "—")
            add_text(s3, name,    0.35, yy, 2.8, row_h, size=8.5, bold=True, color=NAVY)
            add_text(s3, family,  3.2,  yy, 2.6, row_h, size=8,   color=DKTXT)
            add_text(s3, cstatus, 5.85, yy, 2.2, row_h, size=8,   color=DKTXT)
            add_text(s3, ref,     8.1,  yy, 1.5, row_h, size=7.5, color=RGBColor(0x64,0x74,0x8B))

        if len(open_cases) > 13:
            add_text(s3, f"… and {len(open_cases) - 13} more cases not shown", 0.35, 1.15 + 13 * row_h, 9, 0.28, size=8, color=RGBColor(0x64,0x74,0x8B))

        slide_footer(s3)

        # ── Slide 4: Conclusion ──────────────────────────────────────────────────
        s4 = prs.slides.add_slide(blank)
        add_rect(s4, 0, 0, 10, 5.625, NAVY)
        add_rect(s4, 0, 2.6, 10, 0.06, GOLD)
        add_text(s4, "SUMMARY", 1, 0.5, 8, 0.5, size=13, bold=True, color=GOLD, align="center")
        add_text(s4, f"Total registered cases: {len(cases)}", 0.8, 1.1, 8.4, 0.38, size=12, color=WHITE)
        add_text(s4, f"Open / active cases: {len(open_cases)}", 0.8, 1.52, 8.4, 0.38, size=12, color=WHITE)
        add_text(s4, f"Cases with overdue SLA: {len(overdue)}", 0.8, 1.94, 8.4, 0.38, size=12, color=RGBColor(0xFC,0xA5,0xA5))
        add_text(s4, f"Senior executive cases: {len(senior)}", 0.8, 2.36, 8.4, 0.38, size=12, color=WHITE)
        add_text(s4, "Way Forward", 0.8, 2.82, 8.4, 0.34, size=12, bold=True, color=WHITE)
        add_text(s4, "• Strengthen OPSC turnaround time for resolving disciplinary cases\n"
                     "• Address PSDB member availability and 28-day notice requirements\n"
                     "• Implement more rigorous disciplinary case tracking systems\n"
                     "• Promote ethics and professionalism across all ministries",
                 0.9, 3.2, 8.2, 1.8, size=11, color=RGBColor(0xCA,0xDC,0xFC))
        add_text(s4, f"Public Service Commission — Compliance & Integrity Unit  |  {date.today().strftime('%d %B %Y')}",
                 0, 5.35, 10, 0.25, size=8, color=RGBColor(0x7F,0xA8,0xC9), align="center")

        # ── Stream response ──────────────────────────────────────────────────────
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        fname = f"PSC_Compliance_Cases_{date.today().strftime('%Y%m%d')}.pptx"
        response = HttpResponse(buf.read(), content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        response["Content-Disposition"] = f'attachment; filename="{fname}"'
        return response

    @action(detail=False, methods=["get"], url_path="analytics")
    def analytics(self, request):
        """FR-10: Live compliance analytics — volumes, stage timing, overdue by role, trend."""
        from collections import defaultdict
        from datetime import timedelta

        from django.db.models import Avg, Count, ExpressionWrapper, F, FloatField, Q
        from django.db.models.functions import TruncMonth
        from django.utils import timezone

        from .compliance_models import (
            CaseFamily,
            ComplianceCaseDecision,
            ComplianceCaseStage,
            ComplianceCaseStatus,
            StageStatus,
        )

        cases_qs = compliance_case_queryset(
            request.user,
            ComplianceCase.objects.select_related("submission"),
        )

        # ── 1. Case volumes by family ────────────────────────────────────────
        family_rows = (
            cases_qs.values("case_family")
            .annotate(total=Count("id"), closed=Count("id", filter=Q(status=ComplianceCaseStatus.CLOSED)))
            .order_by("-total")
        )
        family_labels = {f.value: f.label for f in CaseFamily}
        by_family = [
            {
                "family":       r["case_family"],
                "label":        family_labels.get(r["case_family"], r["case_family"]),
                "total":        r["total"],
                "closed":       r["closed"],
                "active":       r["total"] - r["closed"],
            }
            for r in family_rows
        ]

        # ── 2. Outcomes by type ──────────────────────────────────────────────
        from .compliance_models import ComplianceDecisionOutcome
        outcome_rows = (
            ComplianceCaseDecision.objects
            .filter(case__in=cases_qs)
            .values("outcome")
            .annotate(count=Count("id"))
            .order_by("-count")
        )
        outcome_labels = {o.value: o.label for o in ComplianceDecisionOutcome}
        by_outcome = [
            {"outcome": r["outcome"], "label": outcome_labels.get(r["outcome"], r["outcome"]), "count": r["count"]}
            for r in outcome_rows
        ]

        # ── 3. Average stage completion time (days) by stage code ────────────
        timed_stages = ComplianceCaseStage.objects.filter(
            case__in=cases_qs,
            status=StageStatus.COMPLETED,
            started_at__isnull=False,
            completed_at__isnull=False,
        )
        stage_timing_raw = defaultdict(list)
        for s in timed_stages.values("stage_code", "stage_name", "started_at", "completed_at"):
            delta = (s["completed_at"] - s["started_at"]).days
            if delta >= 0:
                stage_timing_raw[(s["stage_code"], s["stage_name"])].append(delta)

        avg_stage_time = sorted(
            [
                {
                    "stage_code": code,
                    "stage_name": name,
                    "avg_days":   round(sum(vals) / len(vals), 1),
                    "sample_n":   len(vals),
                }
                for (code, name), vals in stage_timing_raw.items()
            ],
            key=lambda x: -x["avg_days"],
        )

        # ── 4. Overdue stages by responsible role ────────────────────────────
        overdue_q = Q(
            case__in=cases_qs,
            sla_status="overdue",
        ) & ~Q(status=StageStatus.COMPLETED)
        overdue_rows = (
            ComplianceCaseStage.objects
            .filter(overdue_q)
            .values("responsible_role")
            .annotate(overdue_count=Count("id"))
            .order_by("-overdue_count")
        )
        overdue_by_role = [
            {
                "role":          r["responsible_role"] or "unassigned",
                "label":         (r["responsible_role"] or "Unassigned").replace("_", " ").title(),
                "overdue_count": r["overdue_count"],
            }
            for r in overdue_rows
        ]

        # Also surface the individual overdue cases for manager attention
        overdue_cases = []
        seen_case_ids = set()
        for s in (
            ComplianceCaseStage.objects
            .filter(overdue_q)
            .select_related("case__submission")
            .order_by("-case__created_at")
        ):
            if s.case_id in seen_case_ids:
                continue
            seen_case_ids.add(s.case_id)
            overdue_cases.append({
                "case_id":        s.case.id,
                "reference":      s.case.submission.reference_number,
                "subject":        s.case.subject_name,
                "ministry":       s.case.subject_ministry,
                "stage_name":     s.stage_name,
                "responsible":    (s.responsible_role or "").replace("_", " ").title(),
                "due_date":       s.due_date.isoformat() if s.due_date else None,
            })

        # ── 5. Caseload trend (last 18 months) ───────────────────────────────
        cutoff = timezone.now() - timedelta(days=548)  # ~18 months
        opened_by_month = (
            cases_qs
            .filter(created_at__gte=cutoff)
            .annotate(month=TruncMonth("created_at"))
            .values("month")
            .annotate(opened=Count("id"))
            .order_by("month")
        )
        closed_by_month = (
            cases_qs
            .filter(date_closed__gte=cutoff)
            .annotate(month=TruncMonth("date_closed"))
            .values("month")
            .annotate(closed=Count("id"))
            .order_by("month")
        )
        trend_map = {}
        for r in opened_by_month:
            k = r["month"].strftime("%Y-%m")
            trend_map.setdefault(k, {"month": k, "opened": 0, "closed": 0})["opened"] = r["opened"]
        for r in closed_by_month:
            k = r["month"].strftime("%Y-%m")
            trend_map.setdefault(k, {"month": k, "opened": 0, "closed": 0})["closed"] = r["closed"]
        caseload_trend = sorted(trend_map.values(), key=lambda x: x["month"])

        # ── Summary totals ────────────────────────────────────────────────────
        total        = cases_qs.count()
        active_count = cases_qs.exclude(status__in=[ComplianceCaseStatus.CLOSED, ComplianceCaseStatus.ARCHIVED]).count()
        closed_count = cases_qs.filter(status=ComplianceCaseStatus.CLOSED).count()
        overdue_count = (
            ComplianceCaseStage.objects
            .filter(overdue_q)
            .values("case_id").distinct().count()
        )

        return Response({
            "summary":       {"total": total, "active": active_count, "closed": closed_count, "overdue_cases": overdue_count},
            "by_family":     by_family,
            "by_outcome":    by_outcome,
            "avg_stage_time": avg_stage_time,
            "overdue_by_role": overdue_by_role,
            "overdue_cases": overdue_cases[:20],
            "caseload_trend": caseload_trend,
        })

    @action(detail=True, methods=["post"], url_path="stage")
    def update_stage(self, request, pk=None):
        """Update a statutory stage's status (e.g. mark complete) and refresh SLA."""
        from django.utils import timezone

        from .compliance_models import StageStatus
        from .compliance_workflows import recompute_case_sla

        case = self.get_object()
        stage_id = request.data.get("stage_id")
        new_status = request.data.get("status")
        stage = case.stages.filter(pk=stage_id).first()
        if not stage:
            raise ValidationError({"stage_id": "Stage not found for this case."})
        if new_status not in StageStatus.values:
            raise ValidationError({"status": "Invalid stage status."})
        stage.status = new_status
        if new_status == StageStatus.IN_PROGRESS and not stage.started_at:
            stage.started_at = timezone.now()
        if new_status == StageStatus.COMPLETED:
            stage.completed_at = timezone.now()
        stage.save(update_fields=["status", "started_at", "completed_at"])
        recompute_case_sla(case)
        return Response(ComplianceCaseDetailSerializer(case, context={"request": request}).data)

    @action(detail=True, methods=["post"], url_path="decisions")
    def record_decision(self, request, pk=None):
        """Record a Commission/PSDB/HOD decision and auto-update case status."""
        from django.utils import timezone as tz

        from .compliance_models import (
            ComplianceCaseDecision,
            ComplianceCaseStatus,
            TERMINAL_OUTCOMES,
        )

        if not user_can_record_decision(request.user):
            raise PermissionDenied("Your role does not have permission to record decisions.")

        case = self.get_object()

        ser = ComplianceCaseDecisionSerializer(data=request.data)
        ser.is_valid(raise_exception=True)
        decision = ser.save(case=case, decided_by=request.user)

        # Auto-close case when outcome is terminal
        if decision.outcome in TERMINAL_OUTCOMES:
            case.status = ComplianceCaseStatus.CLOSED
            case.date_closed = tz.now()
            case.save(update_fields=["status", "date_closed"])

        return Response(
            ComplianceCaseDetailSerializer(case, context={"request": request}).data,
            status=201,
        )


class CanLodgeOrViewComplaint(permissions.BasePermission):
    """Ministry staff may lodge + view their own; compliance staff may do everything."""

    def has_permission(self, request, view):
        user = request.user
        if not user or not user.is_authenticated:
            return False
        if user_can_view_compliance(user):
            return True
        # Ministry staff: may create and read (queryset is scoped to their own).
        if view.action in ("create", "list", "retrieve"):
            return _user_role(user) in MINISTRY_LODGE_ROLES
        return False


class ComplaintViewSet(viewsets.ModelViewSet):
    """Complaints Register (compliance) + ministry write-only lodgement."""

    permission_classes = [permissions.IsAuthenticated, CanLodgeOrViewComplaint]
    http_method_names = ["get", "post", "head", "options"]

    def get_serializer_class(self):
        if self.action == "create" and not _is_compliance_staff(self.request.user):
            return ComplaintLodgeSerializer
        return ComplaintSerializer

    def get_queryset(self):
        base = Complaint.objects.select_related("ministry", "compliance_case__submission").order_by("-created_at")
        return complaint_queryset(self.request.user, base)

    def perform_create(self, serializer):
        user = self.request.user
        from .profile_utils import ensure_psc_profile

        ministry_id = None
        try:
            ministry_id = ensure_psc_profile(user).ministry_id
        except Exception:
            pass
        serializer.save(lodged_by=user, ministry_id=ministry_id)

    @action(detail=True, methods=["post"])
    def accept(self, request, pk=None):
        """Compliance triages a complaint and opens a case from it."""
        if not _is_compliance_staff(request.user):
            raise PermissionDenied("Only Compliance staff may triage complaints.")
        complaint = self.get_object()
        if complaint.compliance_case_id:
            raise ValidationError("This complaint has already been converted to a case.")

        ser = ComplianceCaseCreateSerializer(data={
            **request.data,
            "subject_name": request.data.get("subject_name") or complaint.subject_name or complaint.title,
            "subject_position": request.data.get("subject_position") or complaint.subject_position,
            "subject_ministry": request.data.get("subject_ministry") or complaint.subject_ministry,
            "title": request.data.get("title") or complaint.title,
            "description": request.data.get("description") or complaint.description,
        })
        ser.is_valid(raise_exception=True)
        data = dict(ser.validated_data)
        data.pop("complaint_id", None)
        case = create_compliance_case(creator=request.user, complaint=complaint, **data)
        return Response(
            ComplianceCaseDetailSerializer(case, context={"request": request}).data,
            status=status.HTTP_201_CREATED,
        )

    @action(detail=True, methods=["post"])
    def reject(self, request, pk=None):
        """Compliance rejects a complaint with a reason visible to the lodging ministry."""
        if not _is_compliance_staff(request.user):
            raise PermissionDenied("Only Compliance staff may triage complaints.")
        complaint = self.get_object()
        complaint.status = ComplaintStatus.REJECTED
        complaint.closed_reason = request.data.get("reason", "")
        from django.utils import timezone
        complaint.triaged_by = request.user
        complaint.triaged_at = timezone.now()
        complaint.save(update_fields=["status", "closed_reason", "triaged_by", "triaged_at"])
        return Response(ComplaintSerializer(complaint, context={"request": request}).data)


class OffenceTypeViewSet(viewsets.ModelViewSet):
    """The nature-of-offence catalogue. Readable by compliance staff (for the
    case dropdown); writable only by administrators."""

    permission_classes = [permissions.IsAuthenticated]

    def get_serializer_class(self):
        from .compliance_serializers import OffenceTypeSerializer
        return OffenceTypeSerializer

    def get_queryset(self):
        qs = OffenceType.objects.all()
        if self.request.query_params.get("active") == "true":
            qs = qs.filter(active=True)
        return qs

    def _require_admin(self):
        u = self.request.user
        from .models import Role
        allowed = {Role.PSC_ADMIN, Role.COMPLIANCE_MANAGER}
        if not (u.is_staff or u.is_superuser or _user_role(u) in allowed):
            raise PermissionDenied("Only administrators or the Compliance Manager can modify the offence catalogue.")

    def perform_create(self, serializer):
        self._require_admin()
        serializer.save()

    def perform_update(self, serializer):
        self._require_admin()
        serializer.save()

    def perform_destroy(self, instance):
        self._require_admin()
        instance.delete()
